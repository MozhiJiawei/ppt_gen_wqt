import argparse
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(245, 247, 250)
NAVY = RGBColor(26, 43, 79)
BLUE = RGBColor(67, 113, 196)
TEAL = RGBColor(14, 144, 160)
GREEN = RGBColor(63, 161, 93)
CORAL = RGBColor(224, 102, 80)
GOLD = RGBColor(214, 161, 66)
INK = RGBColor(37, 42, 52)
MUTED = RGBColor(96, 106, 121)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(229, 238, 246)
LIGHT = RGBColor(236, 241, 246)
BORDER = RGBColor(220, 226, 233)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill=WHITE, line=BORDER, radius=True):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def set_card_title(shape, title, size=18, color=NAVY):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color


def add_shape_bullets(shape, items, size=15, color=INK):
    tf = shape.text_frame
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = color


def add_top(slide, title, subtitle, page):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.62))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY
    add_text(slide, Inches(0.62), Inches(0.78), Inches(9.6), Inches(0.52), title, 24, NAVY, True)
    add_text(slide, Inches(0.65), Inches(1.22), Inches(10.4), Inches(0.42), subtitle, 12.5, MUTED)
    tag = add_card(slide, Inches(11.95), Inches(0.14), Inches(0.72), Inches(0.28), TEAL, TEAL)
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{page}/7"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_takeaway(slide, text):
    card = add_card(slide, Inches(0.64), Inches(1.72), Inches(12.05), Inches(0.72), PALE, PALE)
    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "核心结论："
    r1.font.name = "Microsoft YaHei"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = NAVY


def add_bullets(slide, left, top, width, items, size=16, bullet_color=TEAL):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, y + Inches(0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.color.rgb = bullet_color
        add_text(slide, left + Inches(0.2), y, width - Inches(0.2), Inches(0.42), item, size, INK, False)
        y += Inches(0.48)


def add_metric_card(slide, left, top, width, title, value, accent):
    card = add_card(slide, left, top, width, Inches(0.95), WHITE, BORDER)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.22), title, 11.5, MUTED, True)
    add_text(slide, left + Inches(0.2), top + Inches(0.34), width - Inches(0.3), Inches(0.34), value, 22, accent, True)
    return card


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, weight=2.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.end_arrowhead = True
    return line


def save_gif_frame(gif_path: Path, output_path: Path):
    with Image.open(gif_path) as img:
        img.seek(0)
        frame = img.convert("RGB")
        frame.save(output_path)


def build_results_chart(out_path: Path):
    methods = [
        "Base",
        "DGM",
        "No open-ended",
        "No self-improve",
        "Greedy",
    ]
    swe = np.array([20.0, 50.0, 23.0, 39.0, 39.7])
    poly = np.array([14.2, 30.7, 14.0, 28.0, 30.0])

    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=200)
    x = np.arange(len(methods))
    width = 0.34
    colors = ["#D7DEE8", "#3F72C4", "#E06650", "#3FA15D", "#D6A142"]
    ax.bar(x - width / 2, swe, width, label="SWE-bench", color=colors[1], edgecolor="none")
    ax.bar(x + width / 2, poly, width, label="Polyglot", color=colors[2], edgecolor="none")
    ax.set_ylim(0, 60)
    ax.set_xticks(x)
    ax.set_xticklabels(methods, fontsize=10)
    ax.set_ylabel("Success rate (%)", fontsize=10)
    ax.grid(axis="y", color="#D9E1EA", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ["top", "right", "left", "bottom"]:
        ax.spines[spine].set_visible(False)
    ax.tick_params(axis="y", labelsize=9, length=0)
    ax.tick_params(axis="x", length=0)
    ax.legend(loc="upper left", frameon=False, fontsize=10)
    for xpos, val in zip(x - width / 2, swe):
        ax.text(xpos, val + 1.0, f"{val:.1f}", ha="center", va="bottom", fontsize=9, color="#233042")
    for xpos, val in zip(x + width / 2, poly):
        ax.text(xpos, val + 1.0, f"{val:.1f}", ha="center", va="bottom", fontsize=9, color="#233042")
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def build_transfer_chart(out_path: Path):
    labels = ["o3-mini", "Claude 3.5", "Claude 3.7"]
    base = np.array([23.0, 20.0, 19.0])
    dgm = np.array([33.0, 50.0, 59.5])

    fig, ax = plt.subplots(figsize=(5.6, 3.8), dpi=200)
    x = np.arange(len(labels))
    width = 0.34
    ax.bar(x - width / 2, base, width, label="Base agent", color="#D7DEE8", edgecolor="none")
    ax.bar(x + width / 2, dgm, width, label="DGM agent", color="#3F72C4", edgecolor="none")
    ax.set_ylim(0, 65)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel("SWE-bench (%)", fontsize=9)
    ax.grid(axis="y", color="#D9E1EA", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ["top", "right", "left", "bottom"]:
        ax.spines[spine].set_visible(False)
    ax.tick_params(axis="y", labelsize=8, length=0)
    ax.tick_params(axis="x", length=0)
    ax.legend(loc="upper left", frameon=False, fontsize=9)
    for xpos, val in zip(x - width / 2, base):
        ax.text(xpos, val + 1.0, f"{val:.1f}", ha="center", va="bottom", fontsize=8, color="#233042")
    for xpos, val in zip(x + width / 2, dgm):
        ax.text(xpos, val + 1.0, f"{val:.1f}", ha="center", va="bottom", fontsize=8, color="#233042")
    fig.patch.set_facecolor("white")
    plt.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def slide_overview(prs, assets: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "Darwin Gödel Machine", "Open-Ended Evolution of Self-Improving Agents", 1)
    add_takeaway(slide, "它把“改自己”变成可执行的编码任务，再用基准测试验证改动是否值得保留。")

    slide.shapes.add_picture(str(assets["overview"]), Inches(0.78), Inches(2.62), width=Inches(6.8), height=Inches(3.92))

    right = add_card(slide, Inches(7.85), Inches(2.62), Inches(4.8), Inches(3.92), WHITE, BORDER)
    right.shadow.inherit = False
    add_text(slide, Inches(8.1), Inches(2.9), Inches(3.8), Inches(0.28), "为什么值得看", 18, NAVY, True)
    add_bullets(
        slide,
        Inches(8.1),
        Inches(3.28),
        Inches(4.1),
        [
            "把 Gödel Machine 的“可证明改进”替换成“经验验证改进”。",
            "不是只爬一条最优路径，而是保留一整棵可回溯的 agent 档案树。",
            "自动长出更细粒度编辑、长上下文压缩、候选排序等能力。",
        ],
        size=15.5,
    )
    add_metric_card(slide, Inches(8.1), Inches(4.95), Inches(2.0), "SWE-bench", "20.0% → 50.0%", BLUE)
    add_metric_card(slide, Inches(10.28), Inches(4.95), Inches(2.0), "Polyglot", "14.2% → 30.7%", CORAL)


def slide_trigger(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "问题设定", "为什么固定架构 agent 很难持续自我改进", 2)
    add_takeaway(slide, "核心矛盾不是“模型不会写代码”，而是“单次改动很难被可靠评估并积累成长期进化”。")

    card1 = add_card(slide, Inches(0.74), Inches(2.62), Inches(4.2), Inches(3.82), WHITE, BORDER)
    set_card_title(card1, "传统 Gödel Machine", 18)
    add_shape_bullets(card1, [
        "要求证明一次代码重写“净收益为正”。",
        "对真实 LLM agent 来说，证明几乎不可行。",
        "工具收益高度依赖模型、任务与上下文。",
    ], size=15)

    card2 = add_card(slide, Inches(4.56), Inches(2.62), Inches(4.2), Inches(3.82), WHITE, BORDER)
    set_card_title(card2, "普通 agent 优化", 18)
    add_shape_bullets(card2, [
        "常沿着当前最好解一路 hill-climb。",
        "一次坏改动会拖累后续迭代。",
        "被局部最优困住时，很难回头找“踏脚石”。",
    ], size=15)

    card3 = add_card(slide, Inches(8.38), Inches(2.62), Inches(4.24), Inches(3.82), WHITE, BORDER)
    set_card_title(card3, "DGM 的替代方案", 18)
    add_shape_bullets(card3, [
        "把“自我改进”定义成修改自身代码仓库的编码任务。",
        "用基准得分做经验验证，而不是形式证明。",
        "保留完整 archive，让较差节点也可能孕育后续突破。",
    ], size=15)


def slide_mechanism_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "机制一：自指改写闭环", "Self-improve 不是提示词花活，而是对 agent 自身 repo 的真实修改", 3)
    add_takeaway(slide, "同一个 agent 既要会改自己的代码，也要会在外部任务上交付结果，这两个能力被同一个 benchmark 一起约束。")

    steps = [
        ("1 读取日志", "分析当前评测失败模式"),
        ("2 提出功能", "生成下一步自我改进任务"),
        ("3 修改自身", "直接改自己的工具 / workflow / prompt"),
        ("4 重新评测", "在 SWE-bench 或 Polyglot 上跑分"),
        ("5 决定归档", "能编译且保留编辑能力才进入 archive"),
    ]
    lefts = [0.86, 3.36, 5.86, 8.36, 10.86]
    colors = [BLUE, TEAL, GREEN, GOLD, CORAL]
    for idx, ((title, body), left, color) in enumerate(zip(steps, lefts, colors)):
        add_card(slide, Inches(left), Inches(3.05), Inches(1.9), Inches(2.48), WHITE, BORDER)
        pill = add_card(slide, Inches(left + 0.18), Inches(3.22), Inches(1.12), Inches(0.36), color, color)
        tf = pill.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(slide, Inches(left + 0.18), Inches(3.78), Inches(1.48), Inches(0.85), body, 15.5, INK, True)
        if idx < len(steps) - 1:
            add_arrow(slide, Inches(left + 1.9), Inches(4.29), Inches(left + 2.38), Inches(4.29), color=BLUE)

    add_card(slide, Inches(1.08), Inches(5.78), Inches(11.36), Inches(0.52), LIGHT, LIGHT)
    add_text(
        slide,
        Inches(1.3),
        Inches(5.92),
        Inches(10.9),
        Inches(0.2),
        "论文强调：只有“可继续改代码”的 agent 才能活进 archive，因此自我改进能力本身也是被选择出来的。",
        14,
        NAVY,
        True,
        align=PP_ALIGN.CENTER,
    )


def slide_mechanism_archive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "机制二：Archive 不是缓存，是搜索策略", "开放式探索让低分节点也能成为未来高分节点的前驱", 4)
    add_takeaway(slide, "DGM 的关键不是每一步都涨分，而是让所有节点保留非零被选中概率，从而持续探索“有潜力但暂时不优”的支线。")

    canvas = add_card(slide, Inches(0.78), Inches(2.62), Inches(7.05), Inches(3.85), WHITE, BORDER)
    canvas.shadow.inherit = False
    set_card_title(canvas, "示意：一棵会保留分支的 agent 进化树", 16.2)
    positions = {
        "0": (1.35, 4.2),
        "6": (2.7, 3.35),
        "12": (2.7, 5.05),
        "24": (4.35, 3.0),
        "31": (4.35, 4.2),
        "44": (5.95, 2.75),
        "56": (5.95, 4.15),
        "79": (7.0, 3.35),
    }
    edges = [("0", "6"), ("0", "12"), ("6", "24"), ("12", "31"), ("24", "44"), ("31", "56"), ("44", "79")]
    for a, b in edges:
        x1, y1 = positions[a]
        x2, y2 = positions[b]
        add_arrow(slide, Inches(x1), Inches(y1), Inches(x2), Inches(y2), color=RGBColor(156, 171, 190), weight=1.8)

    node_specs = {
        "0": ("20%", RGBColor(198, 211, 228)),
        "6": ("26%", RGBColor(169, 205, 231)),
        "12": ("23%", RGBColor(226, 209, 164)),
        "24": ("31%", RGBColor(118, 190, 157)),
        "31": ("28%", RGBColor(215, 200, 151)),
        "44": ("38%", RGBColor(109, 183, 112)),
        "56": ("34%", RGBColor(118, 190, 157)),
        "79": ("50%", RGBColor(63, 114, 196)),
    }
    for key, (score, fill) in node_specs.items():
        x, y = positions[key]
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.28), Inches(y - 0.28), Inches(0.56), Inches(0.56))
        node.fill.solid()
        node.fill.fore_color.rgb = fill
        node.line.color.rgb = WHITE
        node.line.width = Pt(2.5)
        add_text(slide, Inches(x - 0.18), Inches(y - 0.1), Inches(0.36), Inches(0.16), key, 11.5, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, Inches(x - 0.26), Inches(y + 0.33), Inches(0.52), Inches(0.18), score, 10.5, MUTED, True, PP_ALIGN.CENTER)

    explain = add_card(slide, Inches(8.08), Inches(2.62), Inches(4.55), Inches(3.85), WHITE, BORDER)
    set_card_title(explain, "为什么它比贪心强", 18)
    add_shape_bullets(explain, [
        "父节点选择大致与性能正相关，但与“已有孩子数”负相关。",
        "所有节点都有非零选择概率，旧方案不会被过早丢弃。",
        "论文里第 24 节点之后出现创新爆发，说明踏脚石价值很高。",
        "Ablation 里 DGM Greedy 仅到 39.7%，明显低于 DGM 的 50.0%。",
    ], size=14.8)


def slide_mechanism_capabilities(prs, assets: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "机制三：它到底长出了什么能力", "论文 Figure 3 给出的代表性自我改进几乎都落在工具与工作流层", 5)
    add_takeaway(slide, "DGM 学到的不是单一技巧，而是一整套让 coding agent 更稳、更细、更会反思的工作方式。")

    slide.shapes.add_picture(str(assets["transfer"]), Inches(0.84), Inches(2.9), width=Inches(4.0), height=Inches(2.65))
    add_text(slide, Inches(0.95), Inches(2.66), Inches(3.5), Inches(0.22), "跨模型迁移也保留增益", 16.5, NAVY, True)

    labels = [
        ("更细粒度查看", "按行查看文件，减少整文件阅读噪声", BLUE),
        ("更精确编辑", "字符串替换代替整文件覆写", TEAL),
        ("上下文自压缩", "接近 context limit 时自动总结", GREEN),
        ("多候选补丁", "同题生成多版 patch 再排序", GOLD),
        ("历史感知重试", "后续尝试会参考前一次失败", CORAL),
        ("外部评审", "让另一模型帮助评估和选最优解", BLUE),
    ]
    start_x = 5.2
    start_y = 2.7
    card_w = 2.32
    card_h = 1.12
    gap_x = 0.22
    gap_y = 0.24
    for idx, (title, body, color) in enumerate(labels):
        row = idx // 2
        col = idx % 2
        left = Inches(start_x + col * (card_w + gap_x))
        top = Inches(start_y + row * (card_h + gap_y))
        add_card(slide, left, top, Inches(card_w), Inches(card_h), WHITE, BORDER)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.1), Inches(card_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        add_text(slide, left + Inches(0.18), top + Inches(0.12), Inches(1.75), Inches(0.2), title, 14.5, NAVY, True)
        add_text(slide, left + Inches(0.18), top + Inches(0.43), Inches(1.9), Inches(0.44), body, 11.8, MUTED, False)


def slide_results(prs, assets: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "评测与消融", "主结果、机制拆解与稳定性一起看，结论会更扎实", 6)
    add_takeaway(slide, "DGM 的收益不是偶然涨分；去掉 open-ended exploration 或 self-improve，两条腿都明显变短。")

    slide.shapes.add_picture(str(assets["results"]), Inches(0.8), Inches(2.7), width=Inches(7.2), height=Inches(3.9))

    card_a = add_card(slide, Inches(8.25), Inches(2.7), Inches(4.35), Inches(1.7), WHITE, BORDER)
    set_card_title(card_a, "结果怎么读", 17)
    add_shape_bullets(card_a, [
        "完整 DGM: 50.0% / 30.7%。",
        "无 archive 仅 23.0% / 14.0%，几乎回到起点。",
        "无自我改进时 39.0% / 28.0%，说明 meta-agent 也要进化。",
    ], size=14.2)

    card_b = add_card(slide, Inches(8.25), Inches(4.58), Inches(4.35), Inches(1.82), WHITE, BORDER)
    set_card_title(card_b, "额外证据", 17)
    add_shape_bullets(card_b, [
        "生成 agent 中，具备基础代码编辑能力的比例最高: 51.3%。",
        "Polyglot 三次运行均值 40.7%，标准差 2.3%。",
        "跨模型迁移到 Claude 3.7 Sonnet 仍可到 59.5%。",
    ], size=13.9)


def slide_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "结论与启发", "把论文落回 agent engineering，会得到四个很实际的判断", 7)
    add_takeaway(slide, "真正可持续的自我改进系统，必须同时具备可验证性、可回溯性、可迁移性与基本安全约束。")

    add_card(slide, Inches(0.8), Inches(2.75), Inches(4.05), Inches(3.45), WHITE, BORDER)
    add_text(slide, Inches(1.05), Inches(3.02), Inches(3.3), Inches(0.24), "技术定位", 18, NAVY, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(3.4),
        Inches(3.45),
        [
            "它不是自动调 prompt 的小修小补，而是 agent 代码仓库级别的迭代。",
            "经验验证替代形式证明，是从理论 Gödel Machine 走向工程可行性的关键一步。",
            "开放式搜索让“暂时较差的方案”继续参与未来创新。",
        ],
        size=15,
    )

    right_items = [
        "好的自我改进系统，要把“继续能改自己”也当作筛选条件。",
        "Archive 设计比单点最优更重要，它决定了系统是否能穿越性能低谷。",
        "改进如果能跨模型、跨 benchmark 迁移，才更像真实能力而非 benchmark 过拟合。",
        "安全措施仍然朴素但必要：沙箱、人类监督、只保留可继续编辑代码的 agent。",
    ]
    y = 2.9
    accents = [BLUE, TEAL, GREEN, CORAL]
    for idx, (item, accent) in enumerate(zip(right_items, accents), start=1):
        add_card(slide, Inches(5.15), Inches(y), Inches(7.15), Inches(0.66), WHITE, BORDER)
        badge = add_card(slide, Inches(5.35), Inches(y + 0.12), Inches(0.34), Inches(0.34), accent, accent)
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(idx)
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(slide, Inches(5.83), Inches(y + 0.13), Inches(6.15), Inches(0.22), item, 13.5, INK, True)
        y += 0.82


def build_ppt(output: Path, work_dir: Path):
    asset_dir = work_dir / "_dgm_ppt_assets"
    asset_dir.mkdir(parents=True, exist_ok=True)

    overview_gif = work_dir / "dgm-main" / "misc" / "overview.gif"
    overview_png = asset_dir / "overview_frame.png"
    results_png = asset_dir / "results_chart.png"
    transfer_png = asset_dir / "transfer_chart.png"

    save_gif_frame(overview_gif, overview_png)
    build_results_chart(results_png)
    build_transfer_chart(transfer_png)

    assets = {
        "overview": overview_png,
        "results": results_png,
        "transfer": transfer_png,
    }

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_overview(prs, assets)
    slide_trigger(prs)
    slide_mechanism_loop(prs)
    slide_mechanism_archive(prs)
    slide_mechanism_capabilities(prs, assets)
    slide_results(prs, assets)
    slide_takeaways(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print("PPT generated")


def main():
    parser = argparse.ArgumentParser(description="Generate a DGM technical report PPT.")
    parser.add_argument("output", help="Path to output pptx")
    parser.add_argument(
        "--source-dir",
        default=r"D:\Git_Repo\ppt_gen_wqt\tmp-artifacts\Darwin Gödel Machine",
        help="Directory containing the paper PDF and dgm-main repo",
    )
    args = parser.parse_args()

    build_ppt(Path(args.output), Path(args.source_dir))


if __name__ == "__main__":
    main()
