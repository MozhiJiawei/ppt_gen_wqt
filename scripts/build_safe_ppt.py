import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(245, 247, 250)
NAVY = RGBColor(25, 44, 87)
TEAL = RGBColor(0, 140, 149)
CORAL = RGBColor(225, 94, 74)
INK = RGBColor(37, 42, 52)
MUTED = RGBColor(93, 102, 115)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(229, 238, 246)
CARD_BORDER = RGBColor(224, 229, 235)


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill=WHITE, line=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_top(slide, title, subtitle, page):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.62))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY
    add_text(slide, Inches(0.62), Inches(0.8), Inches(9.7), Inches(0.56), title, 24, NAVY, True)
    add_text(slide, Inches(0.65), Inches(1.26), Inches(10.2), Inches(0.42), subtitle, 13, MUTED)
    tag = add_card(slide, Inches(11.95), Inches(0.14), Inches(0.72), Inches(0.28), TEAL, TEAL)
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{page}/7"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_takeaway(slide, text):
    card = add_card(slide, Inches(0.64), Inches(1.72), Inches(12.05), Inches(0.72), PALE, PALE)
    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "核心结论： "
    r1.font.name = "Microsoft YaHei"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = NAVY


def overview_slide(prs, deck_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, deck_title, "Overview slide skeleton with a safe chart + summary layout", 1)
    add_takeaway(slide, "Replace this sentence with one short thesis. Keep it to a single line when possible.")
    add_card(slide, Inches(0.72), Inches(2.62), Inches(7.1), Inches(3.82))
    add_card(slide, Inches(8.0), Inches(2.62), Inches(4.63), Inches(3.82))
    add_text(slide, Inches(0.95), Inches(2.9), Inches(2.6), Inches(0.3), "Visual region", 18, NAVY, True)
    add_text(slide, Inches(1.0), Inches(3.35), Inches(6.45), Inches(2.5),
             "Put the main chart, method figure, or comparison graphic here.\n\n"
             "This region is intentionally large so labels and legends have room.")
    add_text(slide, Inches(8.28), Inches(2.9), Inches(3.0), Inches(0.3), "Support region", 18, NAVY, True)
    add_text(slide, Inches(8.28), Inches(3.35), Inches(3.65), Inches(2.5),
             "- 3 short bullets max\n"
             "- Keep each bullet compact\n"
             "- Avoid paragraph text in this column")


def trigger_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "Trigger Scenario", "Use this template for the problem/background slide", 2)
    add_takeaway(slide, "State why existing methods break before introducing the new method.")
    add_card(slide, Inches(0.72), Inches(2.62), Inches(5.55), Inches(3.82))
    add_card(slide, Inches(6.48), Inches(2.62), Inches(6.15), Inches(3.82))
    add_text(slide, Inches(1.0), Inches(2.92), Inches(3.8), Inches(0.3), "Narrative block", 18, NAVY, True)
    add_text(slide, Inches(1.0), Inches(3.35), Inches(4.9), Inches(2.45),
             "- Explain the real-world setting\n"
             "- Explain the failure mode\n"
             "- Explain why prior optimization signals are insufficient")
    add_text(slide, Inches(6.75), Inches(2.92), Inches(3.8), Inches(0.3), "Comparison visual", 18, NAVY, True)
    add_text(slide, Inches(6.75), Inches(3.35), Inches(5.1), Inches(2.45),
             "Place a chart, feedback-scale comparison, or pipeline diagram here.")


def mechanism_slide(prs, page, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, title, "Mechanism slide skeleton", page)
    add_takeaway(slide, "One mechanism per slide. If the explanation needs more room, split it.")
    add_card(slide, Inches(0.72), Inches(2.62), Inches(12.0), Inches(3.82))
    add_text(slide, Inches(1.0), Inches(2.92), Inches(3.0), Inches(0.3), "Flow / timeline / matrix", 18, NAVY, True)
    add_text(slide, Inches(1.0), Inches(3.35), Inches(10.9), Inches(2.5),
             "Use the center of the slide for the main mechanism visual.\n"
             "Keep labels short.\n"
             "Use one bottom annotation strip instead of multiple long callouts.")


def results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "Evaluation", "Use a multi-panel chart and a short reading guide", 6)
    add_takeaway(slide, "Show only the results that support the thesis of the presentation.")
    add_card(slide, Inches(0.72), Inches(2.62), Inches(8.35), Inches(3.82))
    add_card(slide, Inches(9.27), Inches(2.62), Inches(3.36), Inches(3.82))
    add_text(slide, Inches(1.0), Inches(2.92), Inches(2.4), Inches(0.3), "Result panels", 18, NAVY, True)
    add_text(slide, Inches(1.0), Inches(3.35), Inches(7.6), Inches(2.5),
             "Insert 2-3 aligned charts here.\n"
             "Do not add full-width prose below the chart block.\n"
             "Put interpretation in the right column.")
    add_text(slide, Inches(9.55), Inches(2.92), Inches(2.4), Inches(0.3), "Reading guide", 18, NAVY, True)
    add_text(slide, Inches(9.55), Inches(3.35), Inches(2.7), Inches(2.5),
             "- What improved\n"
             "- Why the result matters\n"
             "- What generalizes")


def takeaway_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top(slide, "Takeaways", "Use cards instead of a dense paragraph", 7)
    add_takeaway(slide, "End with technical essence and engineering implications, not a re-summary of all results.")
    add_card(slide, Inches(0.72), Inches(2.62), Inches(5.95), Inches(3.82))
    add_card(slide, Inches(6.95), Inches(2.62), Inches(5.68), Inches(3.82))
    add_text(slide, Inches(1.0), Inches(2.92), Inches(3.0), Inches(0.3), "Technical positioning", 18, NAVY, True)
    add_text(slide, Inches(1.0), Inches(3.35), Inches(5.0), Inches(2.55),
             "Put the matrix, quadrant, or conceptual diagram here.")
    y = 3.25
    for idx in range(4):
        add_card(slide, Inches(7.25), Inches(y), Inches(5.05), Inches(0.58), WHITE)
        add_text(slide, Inches(7.42), Inches(y + 0.14), Inches(0.28), Inches(0.18), str(idx + 1), 13, CORAL, True, PP_ALIGN.CENTER)
        add_text(slide, Inches(7.82), Inches(y + 0.12), Inches(4.1), Inches(0.18), "Short takeaway text", 13, INK, True)
        y += 0.77


def main():
    parser = argparse.ArgumentParser(description="Generate a safe PPT skeleton for technical reports.")
    parser.add_argument("output", help="Path to output pptx")
    parser.add_argument("--title", default="Technical Report", help="Deck title for the overview slide")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    overview_slide(prs, args.title)
    trigger_slide(prs)
    mechanism_slide(prs, 3, "Core Mechanism 1")
    mechanism_slide(prs, 4, "Core Mechanism 2")
    mechanism_slide(prs, 5, "Core Mechanism 3")
    results_slide(prs)
    takeaway_slide(prs)

    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(output)


if __name__ == "__main__":
    main()
