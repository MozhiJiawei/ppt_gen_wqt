import argparse
import json
from dataclasses import dataclass, asdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400


@dataclass
class Finding:
    slide: int
    shape_index: int
    severity: str
    issue: str
    text: str
    metrics: dict


def emu_to_in(value):
    return round(value / EMU_PER_INCH, 2)


def count_mixed_chars(text):
    score = 0.0
    for ch in text:
        if ch.isspace():
            continue
        if ord(ch) < 128:
            score += 0.6
        else:
            score += 1.0
    return score


def shape_bounds(shape):
    return {
        "left": shape.left,
        "top": shape.top,
        "right": shape.left + shape.width,
        "bottom": shape.top + shape.height,
    }


def intersects(a, b):
    return not (
        a["right"] <= b["left"]
        or a["left"] >= b["right"]
        or a["bottom"] <= b["top"]
        or a["top"] >= b["bottom"]
    )


def outside_margin(inner, outer, tol=0.02):
    tol_emu = int(tol * EMU_PER_INCH)
    return (
        inner["left"] < outer["left"] - tol_emu
        or inner["top"] < outer["top"] - tol_emu
        or inner["right"] > outer["right"] + tol_emu
        or inner["bottom"] > outer["bottom"] + tol_emu
    )


def get_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def mean_font_size_pt(shape, default=18.0):
    if not getattr(shape, "has_text_frame", False):
        return default
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return sum(sizes) / len(sizes) if sizes else default


def estimate_wrapped_lines(text, width_in, font_pt):
    if not text:
        return 0
    weighted_chars = count_mixed_chars(text)
    chars_per_line = max(8.0, (width_in * 13.0) * (18.0 / max(font_pt, 10.0)))
    return int((weighted_chars / chars_per_line) + 0.999)


def classify_box(width_in, height_in):
    if width_in <= 2.4 or height_in <= 0.85:
        return "compact"
    if width_in >= 6.5 and height_in <= 0.7:
        return "banner"
    return "body"


def inspect_ppt(path):
    prs = Presentation(str(path))
    findings = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        overlay_shapes = []
        for idx, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE):
                overlay_shapes.append((idx, shape))

        for idx, shape in overlay_shapes:
            text = get_text(shape)
            if not text:
                continue

            width_in = emu_to_in(shape.width)
            height_in = emu_to_in(shape.height)
            font_pt = mean_font_size_pt(shape)
            lines = text.count("\n") + 1
            est_wrapped = estimate_wrapped_lines(text, width_in, font_pt)
            box_kind = classify_box(width_in, height_in)

            metrics = {
                "width_in": width_in,
                "height_in": height_in,
                "font_pt": round(font_pt, 1),
                "line_breaks": lines,
                "estimated_wrapped_lines": est_wrapped,
                "mixed_chars": round(count_mixed_chars(text), 1),
                "box_kind": box_kind,
            }

            if box_kind == "compact" and est_wrapped > 2:
                findings.append(Finding(slide_idx, idx, "high", "compact_box_overflow_risk", text[:120], metrics))
            elif box_kind == "body" and est_wrapped > 5:
                findings.append(Finding(slide_idx, idx, "medium", "body_box_dense_text_risk", text[:120], metrics))
            elif box_kind == "banner" and est_wrapped > 2:
                findings.append(Finding(slide_idx, idx, "medium", "banner_wrap_risk", text[:120], metrics))

            if height_in <= 0.45 and est_wrapped >= 2:
                findings.append(Finding(slide_idx, idx, "high", "short_box_wrap_risk", text[:120], metrics))

            if width_in <= 1.4 and count_mixed_chars(text) > 10:
                findings.append(Finding(slide_idx, idx, "high", "narrow_box_text_risk", text[:120], metrics))

            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                inner = shape_bounds(shape)
                overlapping_cards = []
                for other_idx, other_shape in overlay_shapes:
                    if other_idx == idx:
                        continue
                    if other_shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                        continue
                    other_text = get_text(other_shape)
                    # Prefer empty decorative cards as likely background containers.
                    if other_text:
                        continue
                    bounds = shape_bounds(other_shape)
                    if intersects(inner, bounds):
                        overlapping_cards.append((other_idx, bounds))
                for other_idx, bounds in overlapping_cards:
                    if outside_margin(inner, bounds, tol=0.01):
                        metrics2 = dict(metrics)
                        metrics2["background_shape_index"] = other_idx
                        findings.append(Finding(slide_idx, idx, "medium", "textbox_exceeds_background_card", text[:120], metrics2))
                        break

    return findings


def main():
    parser = argparse.ArgumentParser(description="Heuristic PPT layout checker for text overflow and overlap risks.")
    parser.add_argument("pptx", help="Path to the PowerPoint file")
    parser.add_argument("--json", action="store_true", help="Emit JSON instead of plain text")
    args = parser.parse_args()

    path = Path(args.pptx)
    findings = inspect_ppt(path)

    if args.json:
        print(json.dumps([asdict(f) for f in findings], ensure_ascii=False, indent=2))
        return

    print(f"PPT: {path}")
    print(f"Findings: {len(findings)}")
    for f in findings:
        print(
            f"[{f.severity}] slide {f.slide} shape {f.shape_index}: {f.issue} | "
            f"{f.text} | metrics={f.metrics}"
        )


if __name__ == "__main__":
    main()
